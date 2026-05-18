from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from pptx_voice_video.audio_utils import audio_duration
from pptx_voice_video.models import SlideRecord
from pptx_voice_video.pointer import plan_slide_pointers
from pptx_voice_video.renderer import create_renderer
from pptx_voice_video.config import RenderConfig
from pptx_voice_video.video_compose import compose_video

ROOT = Path(__file__).resolve().parents[1]
INPUTS = ROOT / "sample_inputs" / "pointer_test_10_smoke"
OUT = ROOT / "sample_outputs" / "pointer_test_10_smoke"
AUDIO = INPUTS / "audio"
INPUTS.mkdir(parents=True, exist_ok=True)
AUDIO.mkdir(parents=True, exist_ok=True)
OUT.mkdir(parents=True, exist_ok=True)
PPTX = INPUTS / "pointer_test_10_smoke.pptx"
NOTES_JSON = INPUTS / "pointer_notes.json"

slide_specs = [
    ("제품 소개 테스트", "오늘은 제품 소개 테스트 제목을 가리킵니다."),
    ("핵심 기능", "여기서는 핵심 기능이라는 텍스트를 보세요."),
    ("처리 흐름", "오른쪽 그림을 보시면 처리 흐름을 이해할 수 있습니다."),
    ("사용 사례", "대표 사용 사례라는 텍스트가 핵심입니다."),
    ("운영 전략", "왼쪽 이미지가 운영 전략을 보여줍니다."),
    ("품질 기준", "품질 기준 텍스트를 보겠습니다."),
    ("결과 분석", "상단 차트를 보면 결과 분석의 흐름이 보입니다."),
    ("주의 사항", "주의 사항이라는 텍스트를 가리키겠습니다."),
    ("명확하지 않은 슬라이드", "이 슬라이드는 전반적인 설명만 하고 명확한 지시가 없습니다."),
    ("정리", "마지막으로 정리 제목을 다시 확인합니다."),
]

prs = Presentation()
blank = prs.slide_layouts[6]
for idx, (title, _note) in enumerate(slide_specs, start=1):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.6))
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(30)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 41, 59)
    if idx == 2:
        slide.shapes.add_textbox(Inches(5.2), Inches(2.4), Inches(3.4), Inches(0.8)).text = "핵심 기능"
    elif idx == 3:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(2.3), Inches(2.4), Inches(1.4))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(96, 165, 250); shape.text = "오른쪽 그림"
    elif idx == 4:
        slide.shapes.add_textbox(Inches(3.2), Inches(3.0), Inches(4.0), Inches(0.8)).text = "대표 사용 사례"
    elif idx == 5:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(2.3), Inches(2.0), Inches(1.5))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(34, 197, 94); shape.text = "왼쪽 이미지"
    elif idx == 6:
        slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(3.5), Inches(0.8)).text = "품질 기준"
    elif idx == 7:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(1.3), Inches(3.8), Inches(1.0))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(251, 191, 36); shape.text = "상단 차트"
    elif idx == 8:
        slide.shapes.add_textbox(Inches(5.0), Inches(3.0), Inches(3.5), Inches(0.8)).text = "주의 사항"
    else:
        slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(7.0), Inches(1.0)).text = "테스트용 본문 영역입니다."
prs.save(PPTX)
NOTES_JSON.write_text(json.dumps({i: note for i, (_title, note) in enumerate(slide_specs, start=1)}, ensure_ascii=False, indent=2), encoding="utf-8")

for idx in range(1, 11):
    wav = AUDIO / f"slide_{idx:03}.wav"
    subprocess.run([
        "ffmpeg", "-hide_banner", "-loglevel", "error", "-y",
        "-f", "lavfi", "-i", f"sine=frequency={420 + idx * 20}:duration=1.2",
        "-ar", "24000", str(wav),
    ], check=True)

renderer = create_renderer("libreoffice", RenderConfig(width=1920, height=1080))
images = renderer.render(PPTX, OUT / "slides")
notes_by_slide = {i: note for i, (_title, note) in enumerate(slide_specs, start=1)}
pointer_plans = plan_slide_pointers(PPTX, notes_by_slide, video_width=1920, video_height=1080)
slides = []
for idx, (title, note) in enumerate(slide_specs, start=1):
    wav = AUDIO / f"slide_{idx:03}.wav"
    plan = pointer_plans[idx]
    slides.append(SlideRecord(
        index=idx,
        title=title,
        raw_notes=note,
        normalized_notes=note,
        image_path=str(images[idx - 1]),
        audio_path=str(wav),
        duration=audio_duration(wav) + 0.35,
        engine_name="test-tone",
        status="synthesized",
        pointer_plan=plan.to_dict(),
    ))
final = compose_video(slides, OUT / "pointer_test_10.mp4", fps=30, width=1920, height=1080, pointer_plans=pointer_plans)
manifest = {"slides": [s.to_dict() for s in slides]}
(OUT / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
report = {"final": str(final), "slides": len(slides), "pointer_enabled": True, "manifest": str(OUT / "manifest.json")}
(OUT / "report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
print(json.dumps(report, ensure_ascii=False, indent=2))
