from __future__ import annotations

from dataclasses import asdict, dataclass
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


@dataclass(frozen=True)
class PointerPlan:
    slide_index: int
    visible: bool
    target_x: int = 0
    target_y: int = 0
    reason: str = "ambiguous"
    matched_text: str = ""
    start_x: int | None = None
    start_y: int | None = None
    move_seconds: float = 0.8
    hide_after_seconds: float | None = None

    def to_dict(self) -> dict[str, object]:
        return asdict(self)


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _norm(text: str) -> str:
    return re.sub(r"[\s\W_]+", "", (text or "").lower(), flags=re.UNICODE)


def _center(shape, slide_width: Emu, slide_height: Emu, video_width: int, video_height: int) -> tuple[int, int]:
    x = (int(shape.left) + int(shape.width) / 2) / int(slide_width) * video_width
    y = (int(shape.top) + int(shape.height) / 2) / int(slide_height) * video_height
    return round(x), round(y)


def _shape_text(shape) -> str:
    if hasattr(shape, "text"):
        return _clean_text(shape.text)
    return ""


def _is_visual_shape(shape) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    text = _shape_text(shape)
    return not text and int(getattr(shape, "width", 0)) > 0 and int(getattr(shape, "height", 0)) > 0


def _interesting_text_shapes(slide) -> list[tuple[object, str]]:
    out: list[tuple[object, str]] = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if text and len(_norm(text)) >= 2:
            out.append((shape, text))
    return out


def _find_best_text_match(slide, note: str) -> tuple[object, str] | None:
    note_norm = _norm(note)
    if not note_norm:
        return None
    candidates: list[tuple[int, int, object, str]] = []
    for shape, text in _interesting_text_shapes(slide):
        text_norm = _norm(text)
        if not text_norm:
            continue
        if text_norm in note_norm:
            # Prefer longer exact matches; tie-breaker by visual order.
            candidates.append((len(text_norm), -int(shape.top), shape, text))
            continue
        # For multi-line/list boxes, allow a line-level match.
        for line in text.splitlines():
            line_norm = _norm(line)
            if len(line_norm) >= 4 and line_norm in note_norm:
                candidates.append((len(line_norm), -int(shape.top), shape, line.strip()))
                break
    if not candidates:
        return None
    _, _, shape, matched = max(candidates, key=lambda item: (item[0], item[1]))
    return shape, matched


def _direction(note: str) -> str | None:
    lowered = note.lower()
    checks = [
        ("right", ["오른쪽", "우측", "right"]),
        ("left", ["왼쪽", "좌측", "left"]),
        ("top", ["위쪽", "상단", "위의", "top"]),
        ("bottom", ["아래쪽", "하단", "아래", "bottom"]),
    ]
    for name, words in checks:
        if any(w in lowered for w in words):
            return name
    return None


def _mentions_visual(note: str) -> bool:
    return any(w in note.lower() for w in ["그림", "사진", "이미지", "도표", "차트", "그래프", "figure", "image", "picture", "chart"])


def _direction_score(shape, direction: str, slide_width: Emu, slide_height: Emu) -> float:
    cx = (int(shape.left) + int(shape.width) / 2) / int(slide_width)
    cy = (int(shape.top) + int(shape.height) / 2) / int(slide_height)
    if direction == "right":
        return cx
    if direction == "left":
        return 1 - cx
    if direction == "top":
        return 1 - cy
    if direction == "bottom":
        return cy
    return 0.0


def _find_directional_visual(slide, note: str, slide_width: Emu, slide_height: Emu) -> object | None:
    direction = _direction(note)
    if not direction or not _mentions_visual(note):
        return None
    visual_shapes = [shape for shape in slide.shapes if _is_visual_shape(shape)]
    if not visual_shapes:
        # Fall back to any shape in the requested area when no picture/object exists.
        visual_shapes = [shape for shape in slide.shapes if int(getattr(shape, "width", 0)) > 0 and int(getattr(shape, "height", 0)) > 0]
    if not visual_shapes:
        return None
    return max(visual_shapes, key=lambda shape: _direction_score(shape, direction, slide_width, slide_height))


def _clamp_point(x: int, y: int, video_width: int, video_height: int) -> tuple[int, int]:
    return max(0, min(video_width - 48, x)), max(0, min(video_height - 48, y))


def _near_target_start(target_x: int, target_y: int, video_width: int, video_height: int) -> tuple[int, int]:
    """Start close to the target, never from an arbitrary screen corner."""
    return _clamp_point(target_x - 90, target_y + 70, video_width, video_height)


def _natural_start(
    previous: tuple[int, int] | None,
    target_x: int,
    target_y: int,
    video_width: int,
    video_height: int,
) -> tuple[int, int]:
    """Use the last cursor position only when the jump looks natural."""
    if previous is None:
        return _near_target_start(target_x, target_y, video_width, video_height)
    px, py = previous
    dx = target_x - px
    dy = target_y - py
    max_jump = min(video_width, video_height) * 0.35
    if (dx * dx + dy * dy) ** 0.5 <= max_jump:
        return _clamp_point(px, py, video_width, video_height)
    return _near_target_start(target_x, target_y, video_width, video_height)


def _ambiguous_hold_plan(
    slide_index: int,
    previous: tuple[int, int] | None,
    *,
    video_width: int,
    video_height: int,
) -> PointerPlan:
    if previous is None:
        return PointerPlan(slide_index=slide_index, visible=False, reason="ambiguous")
    x, y = _clamp_point(previous[0], previous[1], video_width, video_height)
    return PointerPlan(
        slide_index=slide_index,
        visible=True,
        target_x=x,
        target_y=y,
        reason="ambiguous_hold_then_hide",
        start_x=x,
        start_y=y,
        move_seconds=0.001,
        hide_after_seconds=1.2,
    )


def _plan_for_shape(
    slide_index: int,
    shape,
    *,
    reason: str,
    matched_text: str = "",
    slide_width: Emu,
    slide_height: Emu,
    video_width: int,
    video_height: int,
    move_seconds: float,
    start: tuple[int, int] | None = None,
) -> PointerPlan:
    x, y = _center(shape, slide_width, slide_height, video_width, video_height)
    start_x, start_y = start or _near_target_start(x, y, video_width, video_height)
    return PointerPlan(
        slide_index=slide_index,
        visible=True,
        target_x=x,
        target_y=y,
        reason=reason,
        matched_text=matched_text,
        start_x=start_x,
        start_y=start_y,
        move_seconds=move_seconds,
    )


def plan_slide_pointers(
    pptx: Path,
    notes_by_slide: dict[int, str],
    *,
    video_width: int = 1920,
    video_height: int = 1080,
    move_seconds: float = 0.8,
) -> dict[int, PointerPlan]:
    """Create conservative automatic cursor plans for PPTX slides.

    Rules are intentionally simple:
    1. Slide 1 points at the slide title.
    2. If notes explicitly mention visible slide text, point to that text.
    3. If notes explicitly say a direction + visual object (e.g. "오른쪽 그림"), point there.
    4. If unclear, hold the previous pointer position briefly and then hide it.
    """
    prs = Presentation(str(pptx))
    plans: dict[int, PointerPlan] = {}
    previous_visible_target: tuple[int, int] | None = None
    for idx, slide in enumerate(prs.slides, start=1):
        note = notes_by_slide.get(idx, "") or ""
        if idx == 1 and slide.shapes.title is not None:
            plans[idx] = _plan_for_shape(
                idx,
                slide.shapes.title,
                reason="first_slide_title",
                matched_text=_shape_text(slide.shapes.title),
                slide_width=prs.slide_width,
                slide_height=prs.slide_height,
                video_width=video_width,
                video_height=video_height,
                move_seconds=move_seconds,
            )
            previous_visible_target = (plans[idx].target_x, plans[idx].target_y)
            continue
        if idx == 1:
            text_shapes = _interesting_text_shapes(slide)
            if text_shapes:
                shape, text = text_shapes[0]
                plans[idx] = _plan_for_shape(
                    idx,
                    shape,
                    reason="first_slide_title",
                    matched_text=text,
                    slide_width=prs.slide_width,
                    slide_height=prs.slide_height,
                    video_width=video_width,
                    video_height=video_height,
                    move_seconds=move_seconds,
                )
                previous_visible_target = (plans[idx].target_x, plans[idx].target_y)
                continue

        matched = _find_best_text_match(slide, note)
        if matched:
            shape, text = matched
            target_x, target_y = _center(shape, prs.slide_width, prs.slide_height, video_width, video_height)
            plans[idx] = _plan_for_shape(
                idx,
                shape,
                reason="matching_text",
                matched_text=text,
                slide_width=prs.slide_width,
                slide_height=prs.slide_height,
                video_width=video_width,
                video_height=video_height,
                move_seconds=move_seconds,
                start=_natural_start(previous_visible_target, target_x, target_y, video_width, video_height),
            )
            previous_visible_target = (plans[idx].target_x, plans[idx].target_y)
            continue

        visual = _find_directional_visual(slide, note, prs.slide_width, prs.slide_height)
        if visual:
            target_x, target_y = _center(visual, prs.slide_width, prs.slide_height, video_width, video_height)
            plans[idx] = _plan_for_shape(
                idx,
                visual,
                reason="directional_shape",
                slide_width=prs.slide_width,
                slide_height=prs.slide_height,
                video_width=video_width,
                video_height=video_height,
                move_seconds=move_seconds,
                start=_natural_start(previous_visible_target, target_x, target_y, video_width, video_height),
            )
            previous_visible_target = (plans[idx].target_x, plans[idx].target_y)
            continue

        plans[idx] = _ambiguous_hold_plan(
            idx,
            previous_visible_target,
            video_width=video_width,
            video_height=video_height,
        )
    return plans
