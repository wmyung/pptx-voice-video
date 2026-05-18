from __future__ import annotations
import json, re, zipfile
from pathlib import Path
from xml.etree import ElementTree as ET
from pptx import Presentation
from .models import SlideRecord

NS = {"a":"http://schemas.openxmlformats.org/drawingml/2006/main", "p":"http://schemas.openxmlformats.org/presentationml/2006/main", "r":"http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
REL_NS = {"rel":"http://schemas.openxmlformats.org/package/2006/relationships"}

def _text_from_xml(xml: bytes) -> str:
    root=ET.fromstring(xml)
    texts=[n.text or "" for n in root.findall('.//a:t', NS)]
    return "\n".join(t.strip() for t in texts if t and t.strip()).strip()

def _resolve_notes_target(rel_path: str, target: str) -> str:
    if target.startswith('../'):
        return 'ppt/' + target[3:]
    return str(Path(rel_path).parent / target).replace('\\', '/')

def _notes_map(pptx: Path) -> dict[int,str]:
    out={}
    with zipfile.ZipFile(pptx) as z:
        slide_rels=[n for n in z.namelist() if re.match(r"ppt/slides/_rels/slide\d+\.xml.rels$", n)]
        for rel_path in slide_rels:
            m=re.search(r"slide(\d+)\.xml", rel_path)
            if not m:
                continue
            idx=int(m.group(1))
            rel_root=ET.fromstring(z.read(rel_path))
            for rel in rel_root.findall('rel:Relationship', REL_NS):
                if rel.attrib.get('Type','').endswith('/notesSlide'):
                    target=_resolve_notes_target(rel_path, rel.attrib['Target'])
                    if target in z.namelist():
                        out[idx]=_text_from_xml(z.read(target))
    return out

def extract_notes(pptx: Path, output_json: Path | None=None) -> list[SlideRecord]:
    prs=Presentation(str(pptx))
    notes=_notes_map(pptx)
    records=[]
    for i, slide in enumerate(prs.slides, start=1):
        title=""
        if slide.shapes.title and getattr(slide.shapes.title, 'text', ''):
            title=slide.shapes.title.text.strip()
        elif slide.shapes:
            for s in slide.shapes:
                if hasattr(s,'text') and s.text.strip():
                    title=s.text.strip().splitlines()[0]
                    break
        records.append(SlideRecord(index=i, title=title, raw_notes=notes.get(i,""), status="notes_extracted"))
    if output_json:
        output_json.parent.mkdir(parents=True, exist_ok=True)
        output_json.write_text(json.dumps([r.to_dict() for r in records], ensure_ascii=False, indent=2), encoding='utf-8')
    return records
