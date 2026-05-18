from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from pptx_voice_video.pointer import plan_slide_pointers


def _sample_deck(path):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    title = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(0.8))
    title.text = "주제 소개"
    slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1)).text = "본문"

    slide2 = prs.slides.add_slide(blank)
    slide2.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(4), Inches(0.8)).text = "원인"
    slide2.shapes.add_textbox(Inches(5.2), Inches(2), Inches(3), Inches(1)).text = "핵심 기능"

    slide3 = prs.slides.add_slide(blank)
    slide3.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(4), Inches(0.8)).text = "정리"
    slide3.shapes.add_textbox(Inches(5.8), Inches(2.3), Inches(2), Inches(1)).text = "오른쪽 그림"
    slide3.shapes.add_shape(1, Inches(6), Inches(3), Inches(2), Inches(1.2))

    slide4 = prs.slides.add_slide(blank)
    slide4.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(4), Inches(0.8)).text = "애매함"
    slide4.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(1)).text = "일반 설명"

    prs.save(path)


def test_auto_pointer_first_slide_points_to_title(tmp_path):
    deck = tmp_path / "deck.pptx"
    _sample_deck(deck)

    plans = plan_slide_pointers(
        deck,
        {
            1: "오늘은 주제 소개를 시작합니다.",
            2: "핵심 기능을 보세요.",
            3: "오른쪽 그림을 보시면 됩니다.",
            4: "전반적으로 설명합니다.",
        },
        video_width=1920,
        video_height=1080,
    )

    assert plans[1].visible is True
    assert plans[1].reason == "first_slide_title"
    assert 100 < plans[1].target_x < 1200
    assert plans[1].target_y < 250


def test_auto_pointer_matches_text_and_hides_ambiguous_slides(tmp_path):
    deck = tmp_path / "deck.pptx"
    _sample_deck(deck)

    plans = plan_slide_pointers(
        deck,
        {
            1: "소개합니다.",
            2: "여기서 핵심 기능을 보세요.",
            3: "오른쪽 그림을 보시면 됩니다.",
            4: "전반적으로 설명합니다.",
        },
        video_width=1920,
        video_height=1080,
    )

    assert plans[2].visible is True
    assert plans[2].reason == "matching_text"
    assert plans[2].target_x > 900

    assert plans[3].visible is True
    assert plans[3].reason in {"directional_shape", "matching_text"}
    assert plans[3].target_x > 900

    assert plans[4].visible is True
    assert plans[4].reason == "ambiguous_hold_then_hide"
    assert plans[4].start_x == plans[4].target_x
    assert plans[4].start_y == plans[4].target_y
    assert plans[4].hide_after_seconds == 1.2
