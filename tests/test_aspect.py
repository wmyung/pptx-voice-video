from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pptx_voice_video.aspect import dimensions_for_slide_aspect, pptx_aspect_ratio


def _make_deck(path: Path, width, height) -> None:
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path)


def test_pptx_aspect_ratio_reads_widescreen_deck(tmp_path: Path):
    deck = tmp_path / "wide.pptx"
    _make_deck(deck, Inches(13.333333), Inches(7.5))

    ratio = pptx_aspect_ratio(deck)

    assert round(ratio, 2) == 1.78


def test_dimensions_for_slide_aspect_keeps_width_and_derives_even_height(tmp_path: Path):
    deck = tmp_path / "wide.pptx"
    _make_deck(deck, Inches(13.333333), Inches(7.5))

    width, height = dimensions_for_slide_aspect(deck, width=1280, height=1280)

    assert width == 1280
    assert height == 720
    assert height % 2 == 0


def test_dimensions_for_slide_aspect_handles_four_by_three(tmp_path: Path):
    deck = tmp_path / "standard.pptx"
    _make_deck(deck, Inches(10), Inches(7.5))

    width, height = dimensions_for_slide_aspect(deck, width=1440, height=1080)

    assert width == 1440
    assert height == 1080
