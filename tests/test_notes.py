from pathlib import Path
from pptx import Presentation
from pptx_voice_video.notes_extractor import extract_notes

def test_extract_notes_without_notes(tmp_path: Path):
    p=tmp_path/'x.pptx'; prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text='Title'; prs.save(p)
    recs=extract_notes(p)
    assert recs[0].title == 'Title' and recs[0].raw_notes == ''
